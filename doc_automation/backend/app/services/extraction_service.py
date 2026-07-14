"""
services/extraction_service.py
───────────────────────────────
Extracts plain text from any supported source document format.

Supported formats
─────────────────
• .pdf   → pdfplumber (text + reading-order reconstruction)
• .docx  → python-docx (paragraphs + tables)
• .pptx  → python-pptx (all slides, shapes, tables, notes)
• .doc / .ppt → best-effort fallback via pdfplumber (if converted) or raw bytes

Public API
──────────
    svc = ExtractionService()
    text: str = svc.extract(Path("report.pdf"))

The returned string is clean, UTF-8, with logical section breaks.
It is intentionally NOT trimmed – the AI service will handle that.
"""

from __future__ import annotations

import io
import re
from pathlib import Path
from typing import Optional

from app.core.logging import get_logger

log = get_logger(__name__)

# ── Optional heavy imports (graceful degradation) ─────────────────────────────
try:
    import pdfplumber  # type: ignore
    _HAVE_PDFPLUMBER = True
except ImportError:
    _HAVE_PDFPLUMBER = False
    log.warning("pdfplumber not installed – PDF extraction degraded.")

try:
    from docx import Document as DocxDocument  # type: ignore
    _HAVE_DOCX = True
except ImportError:
    _HAVE_DOCX = False
    log.warning("python-docx not installed – DOCX extraction unavailable.")

try:
    from pptx import Presentation  # type: ignore
    from pptx.util import Pt       # type: ignore
    _HAVE_PPTX = True
except ImportError:
    _HAVE_PPTX = False
    log.warning("python-pptx not installed – PPTX extraction unavailable.")


# ──────────────────────────────────────────────────────────────────────────────
# Extraction Service
# ──────────────────────────────────────────────────────────────────────────────

class ExtractionService:
    """
    Stateless service – create once and call `.extract()` many times.
    Thread-safe (no shared mutable state).
    """

    # Maximum characters returned per file (prevents token overload in AI step)
    MAX_CHARS: int = 60_000

    def extract(self, path: Path) -> str:
        """
        Auto-detect format and extract text.

        Returns
        -------
        str
            Clean, UTF-8 text with section separators.  Empty string if nothing
            could be extracted.
        """
        ext = path.suffix.lower()
        log.debug("Extracting %s (format=%s)", path.name, ext)

        try:
            if ext == ".pdf":
                text = self._extract_pdf(path)
            elif ext in (".docx", ".doc"):
                text = self._extract_docx(path)
            elif ext in (".pptx", ".ppt"):
                text = self._extract_pptx(path)
            else:
                log.warning("Unsupported format '%s' for %s – skipping.", ext, path.name)
                return ""
        except Exception as exc:
            log.error("Extraction failed for %s: %s", path.name, exc, exc_info=True)
            return ""

        cleaned = self._clean(text)
        if len(cleaned) > self.MAX_CHARS:
            log.debug("Truncating %s from %d → %d chars", path.name, len(cleaned), self.MAX_CHARS)
            cleaned = cleaned[: self.MAX_CHARS]

        log.info("Extracted %d chars from %s", len(cleaned), path.name)
        return cleaned

    # ──────────────────────────────────────────────────────────────────────────
    # PDF extraction (pdfplumber)
    # ──────────────────────────────────────────────────────────────────────────

    def _extract_pdf(self, path: Path) -> str:
        if not _HAVE_PDFPLUMBER:
            raise RuntimeError("pdfplumber is required for PDF extraction.")

        parts: list[str] = []
        with pdfplumber.open(str(path)) as pdf:
            for page_num, page in enumerate(pdf.pages, start=1):
                # ── Attempt structured word extraction first ───────────────
                words = page.extract_words(
                    x_tolerance=3,
                    y_tolerance=3,
                    keep_blank_chars=False,
                    use_text_flow=True,
                )
                if words:
                    page_text = self._words_to_text(words)
                else:
                    # Fallback: naive text extract
                    page_text = page.extract_text(x_tolerance=2, y_tolerance=2) or ""

                # ── Extract tables separately ─────────────────────────────
                table_texts: list[str] = []
                try:
                    for table in page.extract_tables():
                        if table:
                            table_texts.append(self._table_to_text(table))
                except Exception:
                    pass  # non-critical

                combined = page_text
                if table_texts:
                    combined += "\n\n--- TABLE ---\n" + "\n--- TABLE ---\n".join(table_texts)

                if combined.strip():
                    parts.append(f"=== Page {page_num} ===\n{combined}")

        return "\n\n".join(parts)

    def _words_to_text(self, words: list[dict]) -> str:
        """Reconstruct reading-order text from pdfplumber word dicts."""
        if not words:
            return ""
        # Sort by top (y) then left (x) – left-to-right, top-to-bottom
        sorted_words = sorted(words, key=lambda w: (round(w["top"], 0), w["x0"]))
        lines: list[list[str]] = []
        current_line: list[str] = []
        last_top: Optional[float] = None

        for w in sorted_words:
            top = round(w["top"], 0)
            if last_top is None or abs(top - last_top) > 5:
                if current_line:
                    lines.append(current_line)
                current_line = [w["text"]]
            else:
                current_line.append(w["text"])
            last_top = top

        if current_line:
            lines.append(current_line)

        return "\n".join(" ".join(ln) for ln in lines)

    def _table_to_text(self, table: list[list]) -> str:
        """Convert a pdfplumber table (list of rows) to readable text."""
        rows: list[str] = []
        for row in table:
            cells = [str(c).strip() if c is not None else "" for c in row]
            rows.append(" | ".join(cells))
        return "\n".join(rows)

    # ──────────────────────────────────────────────────────────────────────────
    # DOCX extraction (python-docx)
    # ──────────────────────────────────────────────────────────────────────────

    def _extract_docx(self, path: Path) -> str:
        if not _HAVE_DOCX:
            raise RuntimeError("python-docx is required for DOCX extraction.")

        doc    = DocxDocument(str(path))
        parts: list[str] = []

        # ── Paragraphs ────────────────────────────────────────────────────────
        for para in doc.paragraphs:
            text = para.text.strip()
            if not text:
                continue
            # Detect headings
            style = para.style.name if para.style else ""
            if "Heading" in style or para.runs and para.runs[0].bold:
                parts.append(f"\n## {text}\n")
            else:
                parts.append(text)

        # ── Tables ────────────────────────────────────────────────────────────
        for table in doc.tables:
            parts.append("\n--- TABLE ---")
            for row in table.rows:
                cells = [c.text.strip() for c in row.cells]
                parts.append(" | ".join(cells))
            parts.append("--- END TABLE ---\n")

        return "\n".join(parts)

    # ──────────────────────────────────────────────────────────────────────────
    # PPTX extraction (python-pptx)
    # ──────────────────────────────────────────────────────────────────────────

    def _extract_pptx(self, path: Path) -> str:
        if not _HAVE_PPTX:
            raise RuntimeError("python-pptx is required for PPTX extraction.")

        prs   = Presentation(str(path))
        parts: list[str] = []

        for slide_num, slide in enumerate(prs.slides, start=1):
            slide_parts: list[str] = [f"=== Slide {slide_num} ==="]

            for shape in slide.shapes:
                # ── Text frames ───────────────────────────────────────────
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        line = " ".join(run.text for run in para.runs).strip()
                        if line:
                            slide_parts.append(line)

                # ── Tables ────────────────────────────────────────────────
                elif shape.has_table:
                    slide_parts.append("--- TABLE ---")
                    for row in shape.table.rows:
                        cells = [cell.text.strip() for cell in row.cells]
                        slide_parts.append(" | ".join(cells))
                    slide_parts.append("--- END TABLE ---")

            # ── Slide notes ────────────────────────────────────────────────
            if slide.has_notes_slide:
                notes_tf = slide.notes_slide.notes_text_frame
                if notes_tf:
                    note_text = notes_tf.text.strip()
                    if note_text:
                        slide_parts.append(f"[Notes]: {note_text}")

            if len(slide_parts) > 1:   # more than just the header
                parts.append("\n".join(slide_parts))

        return "\n\n".join(parts)

    # ──────────────────────────────────────────────────────────────────────────
    # Text cleaning
    # ──────────────────────────────────────────────────────────────────────────

    def _clean(self, text: str) -> str:
        """
        Normalise whitespace and remove common PDF artefacts:
          • null bytes / non-printable control chars
          • excessive blank lines (>2 consecutive)
          • trailing spaces per line
        """
        if not text:
            return ""
        # Remove null bytes and other control characters except newlines/tabs
        text = re.sub(r"[\x00-\x08\x0b\x0c\x0e-\x1f\x7f]", " ", text)
        # Normalise line endings
        text = text.replace("\r\n", "\n").replace("\r", "\n")
        # Strip trailing spaces per line
        text = "\n".join(ln.rstrip() for ln in text.split("\n"))
        # Collapse >2 consecutive blank lines into 2
        text = re.sub(r"\n{3,}", "\n\n", text)
        return text.strip()
